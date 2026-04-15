import os, json, base64, time, uuid, re, hashlib, filetype, logging, sys, io, tempfile, zipfile, tarfile, mimetypes
import requests
from requests.adapters import HTTPAdapter
from azure.core.pipeline.transport import RequestsTransport
import pyarrow.parquet as pq, pandas as pd, httpx
import fitz  # PyMuPDF
from threading import Lock
from collections import OrderedDict
from concurrent.futures import ThreadPoolExecutor, as_completed
from elasticsearch import Elasticsearch, helpers
from azure.storage.blob import BlobServiceClient
from azure.storage.queue import QueueClient
from openai import OpenAI
from pymongo import MongoClient, ReturnDocument
from datetime import datetime, timezone

# ==========================================
# CONFIGURATION & LOGGING
# ==========================================
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s", stream=sys.stdout)
logger = logging.getLogger("LangScope_Indexer")
logging.getLogger("azure").setLevel(logging.WARNING)
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("elastic_transport").setLevel(logging.WARNING)

AZURE_CONN_STR = os.getenv("AZURE_STORAGE_CONN_STR")
MONGO_URI = os.getenv("MONGO_URI")
mongo_client = MongoClient(MONGO_URI) if MONGO_URI else None
ES_USER, ES_PASS = os.getenv("ELASTIC_USER"), os.getenv("ELASTIC_PASS")
ES_URL = os.getenv("ELASTICSEARCH_URL", "http://localhost:9200").strip()
if ".azurecontainerapps.io" in ES_URL: ES_URL = ES_URL.replace("http://", "https://")

LLM_API_KEY, LLM_MODEL_NAME = os.getenv("LLM_API_KEY"), os.getenv("LLM_MODEL_NAME", "llama-3.3-70b-versatile")
LLM_BASE_URL = os.getenv("LLM_BASE_URL", "https://api.groq.com/openai/v1")

# LINE 35 CHANGE: Unified Single Index
SINGLE_INDEX_NAME = "langscope-unified-data"
CONTAINER = "raw-data"

STRUCTURED_EXTS = ('.parquet', '.csv', '.tsv', '.xlsx', '.xls', '.json', '.jsonl', '.zip', '.tar.gz', '.tgz', '.tar')
STREAMABLE_EXTS = ('.csv', '.tsv', '.json', '.jsonl')

# V24: No retries. Fail fast network logic.
es = Elasticsearch(ES_URL, basic_auth=(ES_USER, ES_PASS) if ES_USER else None, request_timeout=120, max_retries=0)
azure_session = requests.Session()
azure_adapter = HTTPAdapter(pool_connections=50, pool_maxsize=50)
azure_session.mount("https://", azure_adapter)

blob_service = BlobServiceClient.from_connection_string(
    AZURE_CONN_STR, 
    max_single_put_connections=20, 
    max_concurrency=20,
    transport=RequestsTransport(session=azure_session)
)
cc_raw = blob_service.get_container_client(CONTAINER)

class AzureBlobStream(io.RawIOBase): 
    def __init__(self, bc): 
        self.s = bc.download_blob(); self.c = self.s.chunks(); self.b = b""
    def readable(self): return True
    def read(self, sz=-1):
        if sz == -1: r = self.b + self.s.readall(); self.b = b""; return r
        while len(self.b) < sz:
            try: self.b += next(self.c)
            except StopIteration: break
        r, self.b = self.b[:sz], self.b[sz:]; return r
    def readinto(self, b):
        chunk = self.read(len(b))
        if not chunk: return 0
        b[:len(chunk)] = chunk
        return len(chunk)
    def close(self):
        if hasattr(self, 's') and hasattr(self.s, 'close'):
            # Force the Azure downloader to release the TCP socket immediately
            self.s.close()
        super().close()

    
class ThreadSafeLRUCache:
    def __init__(self, maxsize=5000): 
        self.cache = OrderedDict(); self.maxsize = maxsize; self.lock = Lock()
    def get(self, key):
        with self.lock:
            if key not in self.cache: return None
            self.cache.move_to_end(key); return self.cache[key]
    def set(self, key, value):
        with self.lock:
            self.cache[key] = value; self.cache.move_to_end(key)
            if len(self.cache) > self.maxsize: self.cache.popitem(last=False)
    def clear(self):
        with self.lock: self.cache.clear()

UPLOADED_MEDIA_CACHE = ThreadSafeLRUCache(maxsize=5000)
MEDIA_EXECUTOR = ThreadPoolExecutor(max_workers=15)

def identify_media(content, file_path):
    """Dynamically identifies media using magic bytes or system MIME fallback."""
    kind = filetype.guess(content)
    if kind and kind.mime.startswith(('image/', 'audio/', 'video/')):
        return True, kind.mime
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type and mime_type.startswith(('image/', 'audio/', 'video/')):
        return True, mime_type
    return False, None

def ensure_idx(domain):
    physical_idx = f"ls-data-{re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')}"
    universal_limit = 5000 
    
    try:
        if not es.indices.exists(index=physical_idx):
            logger.info(f"[SCHEMA] Creating Primary Index: {physical_idx} | Limit: {universal_limit}")
            es.indices.create(index=physical_idx, body={
                "settings": {
                    "index.refresh_interval": "60s",
                    "index.number_of_replicas": 0,
                    "index.mapping.total_fields.limit": universal_limit
                },
                "mappings": {
                    "dynamic_templates": [
                        {"numeric_discovery": {"match_mapping_type": "double", "mapping": {"type": "double"}}},
                        {"long_discovery": {"match_mapping_type": "long", "mapping": {"type": "long"}}}
                    ],
                    "properties": {
                        "domain": {"type": "keyword"}, 
                        "task_type": {"type": "keyword"},
                        "parent_task_category": {"type": "keyword"},
                        "source_file": {"type": "keyword"}, 
                        "search_content": {"type": "text", "analyzer": "english"}, 
                        "raw_payload": {"type": "object", "dynamic": True} # <--- Primary Math Engine
                    }
                }
            })
            es.indices.put_alias(index=physical_idx, name=SINGLE_INDEX_NAME)
            logger.info(f"[ALIAS] Linked {physical_idx} -> {SINGLE_INDEX_NAME}")
    except Exception as e:
        if "resource_already_exists_exception" not in str(e): logger.error(f"Index creation failed: {e}")
            
    return physical_idx

def ensure_large_idx(large_idx):
    try:
        if not es.indices.exists(index=large_idx):
            logger.info(f"[SCHEMA] Creating Large Data Index: {large_idx} (dynamic: False)")
            es.indices.create(index=large_idx, body={
                "settings": {
                    "index.refresh_interval": "60s",
                    "index.number_of_replicas": 0,
                    "index.mapping.total_fields.limit": 5000
                },
                "mappings": {
                    "properties": {
                        "domain": {"type": "keyword"}, 
                        "task_type": {"type": "keyword"},
                        "parent_task_category": {"type": "keyword"},
                        "source_file": {"type": "keyword"}, 
                        "search_content": {"type": "text", "analyzer": "english"}, 
                        "raw_payload": {"type": "object", "dynamic": False} # <--- Segregated Black Box
                    }
                }
            })
            es.indices.put_alias(index=large_idx, name=SINGLE_INDEX_NAME)
            logger.info(f"[ALIAS] Linked {large_idx} -> {SINGLE_INDEX_NAME}")
    except Exception as e:
        if "resource_already_exists_exception" not in str(e): logger.error(f"Large Index failed: {e}")

def categorize_file(llm, fname, readme, sample, ext):
    f_desc, cats = f"Dynamic File ({ext})", []
    if not llm: return f_desc, cats
    
    has_readme = bool(readme and readme.strip())
    
    try:
        logger.info(f"[GROQ] Generating Domain Categorization via: {fname}")
        
        if has_readme:
            sample_text = sample[:2000]
            prompt = (
                f"You are an expert Data Engineer classifying datasets for a search engine.\n\n"
                f"=== DATASET README ===\n"
                f"{readme.strip()}\n\n"
                f"=== FILE DATA SAMPLE ===\n"
                f"File Name: {fname}\n"
                f"Data Sample: {sample_text}\n\n"
                f"INSTRUCTIONS:\n"
                f"1. Give EQUAL WEIGHT to the README and the Data Sample.\n"
                f"2. Use the README to extract the broad domain, purpose, and overall context of the dataset.\n"
                f"3. Use the Data Sample to extract the exact data format, schema, and specific contents of this file.\n"
                f"4. Provide a brief 'description' that synthesizes BOTH what the overall dataset is about AND what specific data points/format this file contains.\n"
                f"5. Provide a list of relevant 'categories'.\n"
                f"Return strictly in JSON format: {{\"description\": \"...\", \"categories\": [\"...\", \"...\"]}}"
            )
        else:
            sample_text = sample[:10000] 
            prompt = (
                f"You are an expert Data Engineer classifying datasets for a search engine.\n\n"
                f"=== FILE DATA SAMPLE (EXTENDED) ===\n"
                f"File Name: {fname}\n"
                f"Data Sample: {sample_text}\n\n"
                f"INSTRUCTIONS:\n"
                f"1. A README was not provided. You must rely ENTIRELY on this extended Data Sample.\n"
                f"2. Analyze the Data Sample to extract the broad domain, exact data format, schema, and specific contents of this file.\n"
                f"3. Provide a brief 'description' of what this dataset likely contains based on the sample.\n"
                f"4. Provide a list of relevant 'categories'.\n"
                f"Return strictly in JSON format: {{\"description\": \"...\", \"categories\": [\"...\", \"...\"]}}"
            )
        
        res = llm.chat.completions.create(
            model=LLM_MODEL_NAME, 
            response_format={"type": "json_object"}, 
            messages=[{"role": "user", "content": prompt}]
        )
        g_json = json.loads(res.choices[0].message.content)
        f_desc, raw_cats = g_json.get("description", f_desc), g_json.get("categories", [])
        
        def _flatten(item):
            if isinstance(item, dict): return sum([_flatten(v) for v in item.values()], [])
            if isinstance(item, list): return sum([_flatten(v) for v in item], [])
            return [str(item)] if item else []
            
        cats = list(set(_flatten(raw_cats)))
        logger.info(f"Domain Categories Locked: {cats}")
        
    except Exception as e: 
        logger.warning(f"Groq categorization failed, falling back to defaults: {e}")
        
    return f_desc, cats

class Extractor:
    @staticmethod
    def text(fs, ext, max_len=35000):
        t = ""
        try:
            data = fs.read() if hasattr(fs, "read") else (open(fs, "rb").read() if isinstance(fs, str) else fs)
            buf = io.BytesIO(data) if isinstance(data, bytes) else fs
            if ext == '.pdf':
                with fitz.open(stream=data, filetype="pdf") as doc:
                    for p in doc: t += p.get_text() + " "
            elif ext in ('.docx', '.doc'):
                import docx; doc = docx.Document(buf)
                for p in doc.paragraphs: t += p.text + "\n"
            elif ext in ('.pptx', '.ppt'):
                import pptx; ppt = pptx.Presentation(buf)
                for s in ppt.slides:
                    for sh in s.shapes:
                        if hasattr(sh, "text"): t += sh.text + "\n"
            else:
                chk = data[:max_len] if isinstance(data, bytes) else data
                dec = chk.decode('utf-8', errors='replace')
                if dec.count('\x00') < 10: t = dec
        except Exception: pass
        return t[:max_len].strip()

    @staticmethod
    def process_media(k, v, domain):
        raw = v.get('bytes', v) if isinstance(v, dict) else v
        orig_path = v.get('path', '') if isinstance(v, dict) else ''
        
        s_path = re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')
        
        chash = hashlib.md5(raw).hexdigest()
        if (cached := UPLOADED_MEDIA_CACHE.get(chash)): 
            url = cached
        else:
            kind = filetype.guess(raw)
            ext = f'.{kind.extension}' if kind else '.bin'
            bpath = f"raw/{s_path}/media/{uuid.uuid4()}{ext}"
            cc_raw.get_blob_client(bpath).upload_blob(raw, overwrite=True)
            url = f"https://{blob_service.account_name}.blob.core.windows.net/{CONTAINER}/{bpath}"
            UPLOADED_MEDIA_CACHE.set(chash, url)
        
        rich_media_pointer = {
            "file_type": "embedded_parquet_media",
            "file_path": orig_path or url.split('/')[-1],
            "media_url": url,
            "size": len(raw)
        }
        return k, rich_media_pointer

    @staticmethod
    def row(rd, domain):
        c, f = {}, {}
        for k, v in rd.items():
            try:
                if not isinstance(v, (list, dict, tuple)) and pd.isna(v): continue
            except: pass
            
            if isinstance(v, (bytes, bytearray)) or (isinstance(v, dict) and 'bytes' in v): 
                f[MEDIA_EXECUTOR.submit(Extractor.process_media, k, v, domain)] = k
            elif isinstance(v, (dict, list, tuple)):
                try: c[str(k)] = json.dumps(v, default=str)
                except: c[str(k)] = str(v)
            else: 
                if hasattr(v, 'item'):
                    try: v = v.item()
                    except: pass
                c[str(k)] = v if isinstance(v, (bool, int, float)) else str(v)
                
        for res in as_completed(f): 
            k, v = res.result(); c[str(k)] = v
        return c

def get_chunks(fs, ext, domain, csize=100): 
    
    def _parse_logic(stream_target, internal_ext):
        """Your original robust parsing logic preserved exactly."""
        try:
            stream_obj = io.BytesIO(stream_target) if isinstance(stream_target, bytes) else stream_target
            
            if internal_ext.endswith('.parquet'):
                import tempfile, os
                with tempfile.NamedTemporaryFile(delete=False, suffix=".parquet") as tmp:
                    if hasattr(stream_obj, 'read'):
                        while True:
                            chunk = stream_obj.read(8 * 1024 * 1024)
                            if not chunk: break
                            tmp.write(chunk)
                    else:
                        tmp.write(stream_obj)
                    tmp_path = tmp.name
                
                try:
                    for b in pq.ParquetFile(tmp_path).iter_batches(batch_size=csize): 
                        yield b.to_pandas()
                finally:
                    if os.path.exists(tmp_path):
                        os.remove(tmp_path)
                        
            elif internal_ext.endswith(('.csv', '.txt', '.tsv')):
                for c in pd.read_csv(
                    stream_obj, 
                    sep='\t' if internal_ext.endswith('.tsv') else ',', 
                    chunksize=csize, 
                    on_bad_lines='skip',
                    encoding='utf-8',
                    encoding_errors='replace'  # <--- FIX: Forces Pandas to swallow bad bytes instead of crashing
                ): 
                    yield c
                    
            elif internal_ext.endswith(('.xlsx', '.xls')): 
                buf = io.BytesIO(stream_obj.read() if hasattr(stream_obj, 'read') else stream_obj)
                yield pd.read_excel(buf)
                
            elif internal_ext in ('.json', '.jsonl'):
                # ---> THE FIX: Infinite RAM Streaming (No Disk Spooling) <---
                import io, ijson
                
                class PeekableStream(io.RawIOBase):
                    def __init__(self, stream, initial_bytes):
                        self.stream = stream
                        self.buffer = initial_bytes
                    def readable(self): return True
                    def read(self, size=-1):
                        if size == -1:
                            res = self.buffer + (self.stream.read() if hasattr(self.stream, 'read') else b"")
                            self.buffer = b""
                            return res
                        if self.buffer:
                            if size <= len(self.buffer):
                                res, self.buffer = self.buffer[:size], self.buffer[size:]
                                return res
                            else:
                                res = self.buffer
                                self.buffer = b""
                                remainder = size - len(res)
                                return res + (self.stream.read(remainder) if hasattr(self.stream, 'read') else b"")
                        return self.stream.read(size) if hasattr(self.stream, 'read') else b""
                    def readinto(self, b):
                        chunk = self.read(len(b))
                        if not chunk: return 0
                        b[:len(chunk)] = chunk
                        return len(chunk)

                try:
                    if internal_ext == '.jsonl':
                        for c in pd.read_json(stream_obj, lines=True, chunksize=csize): yield c
                    elif internal_ext == '.json':
                        peek_size = 512 * 1024
                        peek_bytes = stream_obj.read(peek_size) if hasattr(stream_obj, 'read') else b""
                        
                        target_prefix = 'item' 
                        is_disguised_jsonl = False
                        
                        try:
                            parser = ijson.parse(io.BytesIO(peek_bytes))
                            for prefix, event, value in parser:
                                if not prefix: 
                                    if event == 'start_array':
                                        target_prefix = 'item'
                                        break
                                    elif event == 'start_map':
                                        is_disguised_jsonl = True
                                        target_prefix = ''
                                        break
                                elif event == 'start_array':
                                    target_prefix = f"{prefix}.item"
                                    break
                        except Exception:
                            pass 
                        
                        # We seamlessly glue the peeked bytes back onto the front of the network stream
                        wrapped_stream = PeekableStream(stream_obj, peek_bytes)
                        logger.info(f"[SCHEMA DISCOVERY] Path: '{target_prefix}' | Disguised JSONL: {is_disguised_jsonl}")

                        if is_disguised_jsonl:
                            try:
                                for c in pd.read_json(wrapped_stream, lines=True, chunksize=csize): yield c
                            except Exception as parse_err:
                                logger.error(f"Disguised JSONL parsing failed: {repr(parse_err)}")
                        else:
                            records = []
                            try:
                                for obj in ijson.items(wrapped_stream, target_prefix):
                                    records.append(obj)
                                    if len(records) >= csize:
                                        df = pd.DataFrame(records)
                                        df.columns = df.columns.astype(str)
                                        yield df
                                        records = []
                                if records:
                                    df = pd.DataFrame(records)
                                    df.columns = df.columns.astype(str)
                                    yield df
                            except Exception as parse_err:
                                logger.error(f"SAX JSON parsing failed: {repr(parse_err)}")
                except Exception as stream_e:
                    logger.error(f"Failed to process stream: {stream_e}")
                                          
        except Exception as e: 
            logger.error(f"Parse failed: {e}")

    # ==========================================
    # MEMORY SAFE ARCHIVE DISPATCHER
    # ==========================================
    
    # CASE 1: ZIP ARCHIVES (Must spool to disk because ZIPs require seeking)
    if ext == '.zip':
        import tempfile, os
        with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp_zip:
            if hasattr(fs, 'read'):
                while True:
                    chunk = fs.read(8 * 1024 * 1024)
                    if not chunk: break
                    tmp_zip.write(chunk)
            else:
                tmp_zip.write(fs)
            zip_path = tmp_zip.name
            
        try:
            with zipfile.ZipFile(zip_path, 'r') as z:
                for member in z.namelist():
                    if z.getinfo(member).is_dir(): continue
                    with z.open(member) as f_obj:
                        inner_ext = os.path.splitext(member)[1].lower()
                        
                        # FIX: Reuse existing global tuple. Stream large data, only peek at unknown files.
                        if inner_ext in STRUCTURED_EXTS:
                            yield from _parse_logic(f_obj, inner_ext)
                        else:
                            peek = f_obj.read(2048)
                            is_media, _ = identify_media(peek, member)
                            
                            if is_media:
                                # THE ROLLING DISK BUFFER (ZIP)
                                import tempfile, os
                                fd, tmp_path = tempfile.mkstemp(suffix=inner_ext)
                                with os.fdopen(fd, 'wb') as f:
                                    f.write(peek)
                                    while True:
                                        chunk = f_obj.read(8 * 1024 * 1024) # Stream 8MB at a time
                                        if not chunk: break
                                        f.write(chunk)
                                # Yield the path to the queue, NOT the heavy bytes
                                yield pd.DataFrame([{"binary_path": tmp_path, "file_name": member}])
                            else:
                                safe_max_bytes = 50 * 1024 * 1024 
                                full_data = peek + f_obj.read(safe_max_bytes)
                                if len(full_data) >= safe_max_bytes:
                                    logger.warning(f"[OOM PREVENTION] Truncated unknown file {member} at 50MB.")
                                yield from _parse_logic(io.BytesIO(full_data), inner_ext)

        finally:
            if os.path.exists(zip_path): os.remove(zip_path)

    # CASE 2: TAR/TAR.GZ (Can stream directly over network without disk!)
    elif ext in ('.tar', '.tar.gz', '.tgz'):
        mode = "r|gz" if 'gz' in ext or 'tgz' in ext else "r|"
        
        stream_src = fs
        if isinstance(fs, bytes): stream_src = io.BytesIO(fs)
        
        with tarfile.open(fileobj=stream_src, mode=mode) as tar:
            for member in tar:
                if not member.isfile(): continue
                f_obj = tar.extractfile(member)
                if f_obj:
                    inner_ext = os.path.splitext(member.name)[1].lower()
                    
                    # FIX: Reuse existing global tuple. Stream large data, only peek at unknown files.
                    if inner_ext in STRUCTURED_EXTS:
                        yield from _parse_logic(f_obj, inner_ext)
                    else:
                        peek = f_obj.read(2048)
                        is_media, _ = identify_media(peek, member.name)
                        
                        if is_media:
                            # THE ROLLING DISK BUFFER (TAR)
                            import tempfile, os
                            fd, tmp_path = tempfile.mkstemp(suffix=inner_ext)
                            with os.fdopen(fd, 'wb') as f:
                                f.write(peek)
                                while True:
                                    chunk = f_obj.read(8 * 1024 * 1024) # Stream 8MB at a time
                                    if not chunk: break
                                    f.write(chunk)
                            # Yield the path to the queue, NOT the heavy bytes
                            yield pd.DataFrame([{"binary_path": tmp_path, "file_name": member.name}])
                        else:
                            full_data = peek + f_obj.read()
                            yield from _parse_logic(io.BytesIO(full_data), inner_ext)

    # CASE 3: STANDARD DATA FILES
    else:
        yield from _parse_logic(fs, ext)


def process_file(bc, fname, domain, idx, f_desc, cats, ext, hf_context="", exec_id="default", is_media_domain=False, task_type="Unknown", parent_category="Unknown"):
    ts, tf = 0, 0
    stream_src = None
    try:
        is_client = hasattr(bc, 'download_blob')
        stream_src = AzureBlobStream(bc) if is_client else bc
        
        if ext in STRUCTURED_EXTS:
            try:
                peek_size = 2 * 1024 * 1024 
                peek_bytes = stream_src.read(peek_size) if hasattr(stream_src, 'read') else b""
                
                if hasattr(stream_src, 'b'): 
                    stream_src.b = peek_bytes + stream_src.b
                elif hasattr(stream_src, 'seek'): 
                    try: stream_src.seek(0)
                    except: pass
                
                if ext in ('.csv', '.tsv', '.jsonl'):
                    rows_in_peek = max(1, peek_bytes.count(b'\n'))
                elif ext == '.json':
                    rows_in_peek = max(1, peek_bytes.count(b'}'))
                else:
                    rows_in_peek = 1000 
                
                est_row_bytes = len(peek_bytes) / rows_in_peek if len(peek_bytes) > 0 else 1024
                safe_csize = max(1, min(1000, int((15 * 1024 * 1024) / est_row_bytes)))
                
            except Exception as e:
                logger.warning(f"Pre-profiler failed: {e}. Defaulting safe load to 5.")
                safe_csize = 5
                
            c_iter = get_chunks(stream_src, ext, domain, csize=safe_csize)
            
            try: first = next(c_iter)
            except StopIteration: return 0, 0
            
            try:
                avg_row_bytes = first.memory_usage(deep=True).sum() / max(1, len(first))
                target_chunk_bytes = 4 * 1024 * 1024 
                safe_chunk_size = max(1, min(1000, int(target_chunk_bytes / avg_row_bytes)))
                logger.info(f"[MEMORY PROFILER] {fname} | Avg Row: {round(avg_row_bytes/(1024*1024), 2)} MB -> Dynamic Chunk Size: {safe_chunk_size}")
                # Estimate nested complexity by counting dictionary keys in the first chunk
                max_complexity = max([str(r).count("':") for r in first.to_dict('records')] + [0])
                
                if max_complexity > 500:
                    safe_chunk_size = min(safe_chunk_size, 15) # Throttle to protect the JVM mapping engine
                    logger.warning(f"[SCHEMA PROFILER] High nesting detected (~{max_complexity} keys). Clamping chunk size to {safe_chunk_size}.")
                    
            except Exception as mem_e:
                logger.warning(f"Memory profiler failed, defaulting to 50: {mem_e}")
                safe_chunk_size = 50
                
            # TRACKERS: Required for the Graceful Fallback / Large Index routing
            in_flight_docs = {} 
            retry_actions = []
            
            def action_gen():
                nonlocal tf
                row_counter = 0  # <--- NEW: Track the exact position in the file
                
                def _yield_df(df):
                    nonlocal tf, row_counter # <--- NEW: Bring counter into scope
                    for r in df.to_dict('records'):
                        row_counter += 1
                        # NEW LINE: Handle binary files extracted via the Rolling Buffer
                        if "binary_path" in r and "file_name" in r:
                            tmp_path = r["binary_path"]
                            try:
                                with open(tmp_path, "rb") as f:
                                    # Process it one at a time
                                    _, p = Extractor.process_media(r["file_name"], {"bytes": f.read(), "path": r["file_name"]}, domain)
                            finally:
                                # VAPORIZE INSTANTLY: Keep the pod's hard drive empty
                                if os.path.exists(tmp_path):
                                    os.remove(tmp_path)
                                    
                        # Fallback just in case standard binary content snuck through
                        elif "binary_content" in r and "file_name" in r:
                            _, p = Extractor.process_media(r["file_name"], {"bytes": r["binary_content"], "path": r["file_name"]}, domain)
                        else:
                            p = Extractor.row(r, domain)
                            
                        doc_payload = None
                        
                        if p.get("file_type") == "direct_hf_media":
                            search_val = p.get("file_path", "")
                            full_search = f"Direct Media File: {search_val}".strip() 
                            doc_id = hashlib.md5((f"{fname}_{row_counter}_{search_val}").encode('utf-8')).hexdigest()
                            
                            doc_payload = {
                                "_index": idx, "_id": doc_id, 
                                "_source": {
                                    "domain": domain, 
                                    "task_type": task_type, 
                                    "parent_task_category": parent_category,
                                    "source_file": fname, 
                                    "search_content": full_search, 
                                    "raw_payload": {**p, "ai_description": f_desc, "ai_categories": cats, "hf_context": hf_context}
                                }
                            }
                        else:
                            # SENIOR FIX: Semantic Anchoring for Needle & Math Reasoning
                            search_parts = [f"This record belongs to the {domain} domain."]
                            p_for_es = {} 
                            has_media = False 
                            
                            for k, val in p.items():
                                # Preserve original numeric types for Numeric Discovery (Math Reasoning)
                                p_for_es[k] = val 
                                
                                if isinstance(val, dict) and val.get("file_type") == "embedded_parquet_media":
                                    search_parts.append(f"Contains media reference: {val.get('file_path', '')}")
                                    p_for_es[k] = json.dumps(val) 
                                    has_media = True
                                elif not str(val).startswith(("http", "hf://")):
                                    # THE ANCHOR: Natural language context for the vectorizer
                                    search_parts.append(f"The {k} is {val}.")
                            
                            row_data_str = " ".join(search_parts).strip()
                            
                            # Logic for determining the final search context
                            if is_media_domain or has_media:
                                full_search = row_data_str
                            else:
                                # Combine the anchors with the AI descriptions and file context
                                full_search = f"{row_data_str} {f_desc} {hf_context}".strip()
                                
                            doc_id = hashlib.md5((f"{fname}_{row_counter}_{row_data_str}").encode('utf-8')).hexdigest()
                            
                            doc_payload = {
                                "_index": idx, "_id": doc_id, 
                                "_source": {
                                    "domain": domain, 
                                    "task_type": task_type, 
                                    "parent_task_category": parent_category,
                                    "source_file": fname, 
                                    "search_content": full_search, 
                                    "raw_payload": {**p_for_es, "ai_description": f_desc, "ai_categories": cats, "hf_context": hf_context}
                                }
                            }
                            
                        try:
                            serialized_bytes = json.dumps(doc_payload, default=str).encode('utf-8')
                            byte_footprint = len(serialized_bytes)
                        except Exception as ser_err:
                            logger.error(f"[PRE-FLIGHT REJECT] Row {doc_id} serialization failed: {ser_err}")
                            tf += 1
                            continue
                            
                        limit_bytes = 25 * 1024 * 1024
                        if byte_footprint > limit_bytes:
                            logger.error(f"[PRE-FLIGHT REJECT] Row {doc_id} skipped! Size: {byte_footprint/(1024*1024):.2f}MB exceeds 25MB limit.")
                            tf += 1
                            continue
                        
                        doc_id = doc_payload["_id"]
                        in_flight_docs[doc_id] = doc_payload # <--- ADD THIS LINE: Lock into RAM tracker
                        
                        yield doc_payload
                
                yield from _yield_df(first)
                for chunk in c_iter: yield from _yield_df(chunk)

            error_details = []
            is_heavy = safe_chunk_size < 101
            active_threads = 1 if is_heavy else 2
            log_interval = 1 if is_heavy else 100000
            
            bulk_generator = helpers.parallel_bulk(
                client=es,
                actions=action_gen(),
                thread_count=active_threads,
                chunk_size=safe_chunk_size,
                queue_size=2 if is_heavy else 4,
                max_chunk_bytes=100 * 1024 * 1024,
                request_timeout=120,
                raise_on_error=False,
                raise_on_exception=False
            )

            for ok, info in bulk_generator:
                op_type = next(iter(info))
                doc_id = info[op_type].get('_id')
                
                if ok: 
                    ts += 1
                    in_flight_docs.pop(doc_id, None) 
                else: 
                    # Extract the structured error object
                    error_obj = info[op_type].get('error', {})
                    err_type = error_obj.get('type', '') if isinstance(error_obj, dict) else ''
                    
                    # Official Elasticsearch schema/mapping exception classes
                    schema_exceptions = (
                        "illegal_argument_exception",       # Thrown for field limit breaches
                        "mapper_parsing_exception",         # Thrown for data type conflicts
                        "strict_dynamic_mapping_exception",  # Thrown if dynamic mapping fails
                        "document_parsing_exception",
                        "unavailable_shards_exception"
                    )
                    
                    # GRACEFUL ROUTING: Catch schema failures via strict type checking
                    if err_type in schema_exceptions:
                        if doc_id in in_flight_docs:
                            failed_doc = in_flight_docs.pop(doc_id)
                            
                            large_idx = f"{idx}-large"
                            failed_doc['_index'] = large_idx
                            
                            retry_actions.append(failed_doc)
                        else:
                            tf += 1
                    else:
                        tf += 1
                        in_flight_docs.pop(doc_id, None)
                        logger.error(f"[FAILED & SKIPPED] Row failed. Type: {err_type} | Reason: {error_obj}")
                        try: error_details.append(json.dumps(info, default=str))
                        except: error_details.append(str(info))
                
                if ts > 0 and ts % log_interval == 0:
                    logger.info(f"[PROGRESS] {fname} indexed {ts:,} rows... (Batch Size: {safe_chunk_size})")
                    
            if retry_actions:
                logger.info(f"[GRACEFUL FALLBACK] Attempting to route {len(retry_actions)} documents to the -large index...")
                large_idx_name = f"{idx}-large"
                ensure_large_idx(large_idx_name)
                
                try:
                    retry_gen = helpers.parallel_bulk(
                        client=es,
                        actions=retry_actions,
                        thread_count=1,
                        chunk_size=100,
                        raise_on_error=False,
                        raise_on_exception=False
                    )
                    
                    for r_ok, r_info in retry_gen:
                        if r_ok:
                            ts += 1
                        else:
                            tf += 1
                            logger.error(f"[LARGE INDEX FALLBACK FAILED] {r_info}")
                except Exception as retry_e:
                    logger.error(f"Failed to execute retry batch: {retry_e}")
                    tf += len(retry_actions)
                    
            if error_details:
                s_path = re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')
                err_path = f"status_signals/{exec_id}/errors/{s_path}/indexer_{fname}_rows.log"
                try: cc_raw.get_blob_client(err_path).upload_blob("\n".join(error_details).encode('utf-8'), overwrite=True)
                except: pass
                    
        else:
            txt = Extractor.text(stream_src, ext, 35000)
            
            s_path = re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')
            bp = f"raw/{s_path}/media/{uuid.uuid4()}{ext or '.bin'}"
            
            if is_client: cc_raw.get_blob_client(bp).start_copy_from_url(bc.url)
            else:
                with open(stream_src, "rb") as f: cc_raw.get_blob_client(bp).upload_blob(f, overwrite=True)
            f_url = f"https://{blob_service.account_name}.blob.core.windows.net/{CONTAINER}/{bp}"
            
            full_search = f"{fname} {txt}".strip()
            
            doc_id = hashlib.md5((fname + txt).encode('utf-8')).hexdigest() 
            doc = {
                "domain": domain, 
                "task_type": task_type, 
                "parent_task_category": parent_category,
                "source_file": fname, 
                "search_content": full_search, 
                "raw_payload": {"file_type": ext.replace('.', ''), "media_url": f_url, "description": f_desc, "categories": cats, "hf_context": hf_context}
            }
            
            try: 
                es.index(index=idx, id=doc_id, document=doc) 
                ts += 1
            except Exception as e: 
                tf += 1
                err_path = f"status_signals/{exec_id}/errors/{s_path}/indexer_{fname}_text.log"
                try: cc_raw.get_blob_client(err_path).upload_blob(str(e).encode('utf-8'), overwrite=True)
                except: pass
            
    except Exception as e: 
        logger.error(f"Worker failed: {e}")
        try:
            s_path = re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')
            safe_exec_id = exec_id if exec_id else "unknown_exec"
            safe_blob_name = re.sub(r'[^a-zA-Z0-9]', '_', fname.lower()).strip('_')
            err_path = f"status_signals/{safe_exec_id}/errors/{s_path}/file_error_{safe_blob_name}_{int(time.time())}.log"
            error_payload = f"FILE PROCESSING ERROR\nDomain: {domain}\nFile: {fname}\nTime: {datetime.now(timezone.utc)}\nError: {str(e)}\n"
            cc_raw.get_blob_client(err_path).upload_blob(error_payload.encode('utf-8'), overwrite=True)
        except Exception as blob_err:
            logger.error(f"Failed to write file error report to Azure: {blob_err}")
            
        finally:
        # Guarantee the TCP socket is returned to the pool the millisecond the file is done
            if isinstance(stream_src, AzureBlobStream):
                try: stream_src.close()
                except: pass

    return ts, tf
 
def run_worker():
    logger.info("[SYSTEM] V51")
    logger.info("=== [LangScope] PARALLEL Indexer Active | Queue Fan-Out Architecture ===")
    
    qc = QueueClient.from_connection_string(AZURE_CONN_STR, queue_name="indexer-trigger-queue")
    llm = OpenAI(base_url=LLM_BASE_URL, api_key=LLM_API_KEY, http_client=httpx.Client(proxies=None, timeout=60.0)) if LLM_API_KEY else None
    
    while True:
        msg_found = False

        for msg in qc.receive_messages(max_messages=1, visibility_timeout=900):
            msg_found = True
            try: payload = json.loads(msg.content if isinstance(msg.content, str) else msg.content.decode('utf-8'))
            except: payload = json.loads(base64.b64decode(msg.content).decode('utf-8'))
            
            raw_target = payload.get("file_to_process", "MASTER_TASK")
            logger.info(f"[QUEUE POP] Pulled message for: {raw_target} | Dequeue Count: {msg.dequeue_count}")
            
            domain, s_path = payload["domain"], payload["s_path"]
            idx = f"ls-data-{re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')}"
            
            task_type = payload.get("task_type", "Unknown")
            parent_category = payload.get("parent_category", "Unknown")

            exec_id = payload.get("exec_id", "default_run")
            sig_prefix = f"status_signals/{exec_id}/indexer/{s_path}"

            if msg.dequeue_count > 2:
                fname = payload.get("file_to_process", "master_task")
                logger.error(f"[POISON EVICTION] {fname} hard-crashed pods {msg.dequeue_count} times. Dropping.")
                
                try: 
                    b_path = fname if fname.startswith("raw/") else f"raw/{s_path}/{fname}"
                    cc_raw.delete_blob(b_path)
                    logger.info(f"[POISON CLEANUP] Vaporized toxic blob: {b_path}")
                except Exception as e:
                    logger.warning(f"[POISON CLEANUP] Skip deletion (already gone or error): {e}")
                    
                try: 
                    err_path = f"status_signals/{exec_id}/errors/{s_path}/poison_evicted_{fname}.log"
                    cc_raw.get_blob_client(err_path).upload_blob(b"Hard pod crash limit exceeded (OOM or Timeout).", overwrite=True)
                    qc.delete_message(msg.id, msg.pop_receipt)
                except: pass
                
                if "file_to_process" in payload:
                    try:
                        if mongo_client:
                            idx_col = mongo_client["langscope"]["indexer_run_history"]
                            # Increment progress ONLY
                            updated_doc = idx_col.find_one_and_update(
                                {"_id": exec_id},
                                {"$inc": {f"domains.{domain}.messages_processed": 1}},
                                return_document=ReturnDocument.AFTER
                            )
                            
                            if updated_doc:
                                domain_data = updated_doc.get("domains", {}).get(domain, {})
                                messages_processed = domain_data.get("messages_processed", 0)
                                expected_count = domain_data.get("expected_files", 0)
                                
                                # Check if this poison pill was the final message
                                if expected_count > 0 and messages_processed >= expected_count:
                                    
                                    # ---> SENIOR FIX: The Atomic Lock <---
                                    claim_lock = idx_col.find_one_and_update(
                                        {
                                            "_id": exec_id, 
                                            f"domains.{domain}.status": {"$nin": ["AGGREGATING", "COMPLETED"]}
                                        },
                                        {"$set": {f"domains.{domain}.status": "AGGREGATING"}},
                                        return_document=ReturnDocument.AFTER
                                    )
                                    
                                    if claim_lock:
                                        logger.info(f"[GRAND TALLY] Triggered by Poison Pill. Lock acquired! Aggregating {domain}...")
                                        total_s, total_f = 0, 0
                                        all_sigs = cc_raw.list_blobs(name_starts_with=sig_prefix, include=['metadata'])
                                        for sig in all_sigs:
                                            if sig.name.endswith(".done") and sig.metadata:
                                                total_s += int(sig.metadata.get("s", 0))
                                                total_f += int(sig.metadata.get("f", 0))

                                        end_time = datetime.now(timezone.utc)
                                        raw_start = updated_doc.get("start_time")
                                        start_time = raw_start.replace(tzinfo=timezone.utc) if raw_start and raw_start.tzinfo is None else (raw_start or end_time)
                                        duration_mins = round((end_time - start_time).total_seconds() / 60, 2)

                                        etl_doc = mongo_client["langscope"]["etl_run_history"].find_one({"_id": exec_id})
                                        total_bytes = etl_doc.get("domains", {}).get(domain, {}).get("total_bytes", 0) if etl_doc else 0

                                        idx_col.update_one(
                                            {"_id": exec_id},
                                            {"$set": {
                                                f"domains.{domain}.rows_success": total_s,
                                                f"domains.{domain}.rows_failed": total_f,
                                                f"domains.{domain}.status": "COMPLETED",
                                                f"domains.{domain}.end_time": end_time,
                                                f"domains.{domain}.duration_minutes": duration_mins,
                                                f"domains.{domain}.total_bytes_processed": total_bytes
                                            }}
                                        )
                                        cc_raw.get_blob_client(f"raw/{s_path}/indexing.done").upload_blob(b"DONE", overwrite=True)
                                        logger.info(f"[FINAL SUMMARY] {domain} | Total Success: {total_s:,} | Total Failed: {total_f:,}")
                    except Exception as e: 
                        logger.error(f"[POISON TALLY CRASH] {e}")
                
                continue

            if "file_to_process" not in payload:
                logger.info(f"[MASTER] Generating Sub-Tasks for Domain: {domain}")
                
                target_idx = f"ls-data-{re.sub(r'[^a-z0-9]', '_', domain.lower()).strip('_')}"
                
                if payload.get("reset"): 
                    try:
                        logger.info(f"[MASTER] Reset requested. Atomic wipe of physical index: {target_idx}")
                        es.indices.delete(index=target_idx, ignore_unavailable=True)
                    except Exception as reset_e:
                        logger.warning(f"Reset failed: {reset_e}")
                
                ensure_idx(domain)
                
                readme = ""
                f_desc = f"Data for {domain}"
                cats = []
                
                hf_context_map = payload.get("hf_context_map", {}) 
                is_media_domain = False 
                
                all_blobs = [b.name for b in cc_raw.list_blobs(name_starts_with=f"raw/{s_path}/") if not any(x in b.name for x in ["indexing.done", "status_signals", ".md", "signature"])]
                if not all_blobs:
                    qc.delete_message(msg.id, msg.pop_receipt); continue
                
                for b in cc_raw.list_blobs(name_starts_with=f"raw/{s_path}/"):
                    if b.name.lower().endswith("readme.md"):
                        readme = cc_raw.get_blob_client(b.name).download_blob().readall().decode('utf-8', errors='ignore'); break

                sample_file = next((b for b in all_blobs if b.endswith(STREAMABLE_EXTS)), all_blobs[0])
                ext = os.path.splitext(sample_file)[1].lower()
                is_media_domain = ext not in STRUCTURED_EXTS
                
                try:
                    if ext in STRUCTURED_EXTS:
                        sample_data = cc_raw.get_blob_client(sample_file).download_blob().readall()
                        first_chunk = next(get_chunks(sample_data, ext, domain, 20))
                        f_desc, cats = categorize_file(llm, sample_file, readme, str(first_chunk.to_dict('records')[:5]), ext)
                
                    if any(c.lower() in ["audio", "video", "image", "media", "multimedia"] for c in cats):
                        is_media_domain = True
                except Exception as e:
                    logger.warning(f"[MASTER] Data profiler skipped or failed for {sample_file}: {e}")
                
                if is_media_domain:
                    logger.info(f"[MASTER] Media Domain detected. Indexing Discovery Header for {domain}")
                    header_id = hashlib.md5((domain + "domain_header").encode('utf-8')).hexdigest()
                    header_search = f"{domain} {f_desc} {' '.join(cats)} {readme[:5000]}".strip()
                    
                    es.index(index=idx, id=header_id, document={
                        "domain": domain, 
                        "task_type": task_type,
                        "parent_task_category": parent_category,
                        "source_file": "DOMAIN_HEADER", 
                        "search_content": header_search,
                        "raw_payload": {
                            "file_type": "header", 
                            "description": f_desc, 
                            "categories": cats, 
                            "full_readme": readme
                        }
                    })

                total_expected = len(all_blobs)
                if mongo_client:
                    mongo_client["langscope"]["indexer_run_history"].update_one(
                        {"_id": exec_id},
                        {
                            "$set": {f"domains.{domain}.expected_files": total_expected},
                            "$setOnInsert": {"start_time": datetime.now(timezone.utc)}
                        },
                        upsert=True
                    )
                cc_raw.get_blob_client(f"{sig_prefix}/expected_files.count").upload_blob(str(total_expected).encode(), overwrite=True)
                
                # FAN OUT: Propagate context maps and taxonomy to workers
                for b in all_blobs:
                    sub_payload = {
                        "domain": domain, 
                        "s_path": s_path, 
                        "file_to_process": b, 
                        "f_desc": f_desc, 
                        "cats": cats, 
                        "hf_context_map": hf_context_map, 
                        "exec_id": exec_id, 
                        "is_media_domain": is_media_domain,
                        "task_type": task_type,
                        "parent_category": parent_category
                    }
                    qc.send_message(base64.b64encode(json.dumps(sub_payload).encode('utf-8')).decode('utf-8'))
                
                logger.info(f"[MASTER] Fanned out {len(all_blobs)} parallel indexing tasks!")
                qc.delete_message(msg.id, msg.pop_receipt)

            else:
                # ==========================================
                # WORKER LOGIC
                # ==========================================
                b_name = payload["file_to_process"]
                f_desc = payload.get("f_desc", "")
                cats = payload.get("cats", [])
                is_media_domain = payload.get("is_media_domain", False) 
                
                task_type = payload.get("task_type", "Unknown")
                parent_category = payload.get("parent_category", "Unknown")
                hf_context_map = payload.get("hf_context_map", {})

                fname, ext = os.path.basename(b_name), os.path.splitext(b_name)[1].lower()
                bc = cc_raw.get_blob_client(b_name)
                
                if not bc.exists():
                    qc.delete_message(msg.id, msg.pop_receipt); continue
                    
                # ---> FIX 3: Dynamic Queue Lease <---
                try:
                    blob_props = bc.get_blob_properties()
                    file_size_gb = blob_props.size / (1024**3)
                    dynamic_timeout = min(int(900 + (file_size_gb * 600)), 14400)
                    qc.update_message(msg.id, msg.pop_receipt, visibility_timeout=dynamic_timeout)
                    logger.info(f"[DYNAMIC LEASE] Extended queue timeout to {dynamic_timeout}s for {file_size_gb:.2f}GB file.")
                except Exception as lease_e:
                    logger.warning(f"Failed to set dynamic queue lease: {lease_e}")
                
                hf_context = ""
                
                hf_context = ""
                for k, v in hf_context_map.items():
                    safe_k = k.replace('/', '_')
                    if fname.startswith(f"hf_{safe_k}") or fname == k:
                        hf_context = v
                        break
                        
                logger.info(f"[WORKER] Indexing: {fname} (Task: {task_type})")
                UPLOADED_MEDIA_CACHE.clear()
                
                try:
                    ts, tf = process_file(bc, fname, domain, idx, f_desc, cats, ext, hf_context, exec_id, is_media_domain, task_type, parent_category)
                    logger.info(f"[FILE DONE] {fname} | Success: {ts:,} | Failed: {tf:,}")
                except Exception as fatal_e:
                    logger.error(f"Fatal Worker Crash on {fname}: {fatal_e}")
                    try:
                        safe_exec_id = exec_id if exec_id else "unknown_exec"
                        err_path = f"status_signals/{safe_exec_id}/errors/{s_path}/fatal_crash_{fname}_{int(time.time())}.log"
                        error_payload = f"FATAL WORKER CRASH\nDomain: {domain}\nFile: {fname}\nTime: {datetime.now(timezone.utc)}\nError: {str(fatal_e)}\n"
                        cc_raw.get_blob_client(err_path).upload_blob(error_payload.encode('utf-8'), overwrite=True)
                    except: pass
                finally:
                    # Existing cleanup
                    try: cc_raw.delete_blob(b_name)
                    except: pass
                    try: qc.delete_message(msg.id, msg.pop_receipt)
                    except: pass
                    
                    # ---> THE FIX: Single Signal with Metadata <---
                    worker_done_id = uuid.uuid4().hex
                    cc_raw.get_blob_client(f"{sig_prefix}/{worker_done_id}.done").upload_blob(
                        b"OK", 
                        overwrite=True,
                        metadata={"s": str(ts), "f": str(tf)} # 's' for success, 'f' for failed
                    )
                    
                    try:
                        if mongo_client:
                            idx_col = mongo_client["langscope"]["indexer_run_history"]
                            # Increment progress ONLY
                            updated_doc = idx_col.find_one_and_update(
                                {"_id": exec_id},
                                {"$inc": {f"domains.{domain}.messages_processed": 1}},
                                return_document=ReturnDocument.AFTER
                            )
                            
                            if updated_doc:
                                domain_data = updated_doc.get("domains", {}).get(domain, {})
                                messages_processed = domain_data.get("messages_processed", 0)
                                expected_count = domain_data.get("expected_files", 0)
                                
                                if expected_count > 0 and messages_processed >= expected_count:
                                    
                                    # ---> SENIOR FIX: The Atomic Lock <---
                                    claim_lock = idx_col.find_one_and_update(
                                        {
                                            "_id": exec_id, 
                                            f"domains.{domain}.status": {"$nin": ["AGGREGATING", "COMPLETED"]}
                                        },
                                        {"$set": {f"domains.{domain}.status": "AGGREGATING"}},
                                        return_document=ReturnDocument.AFTER
                                    )
                                    
                                    if claim_lock:
                                        logger.info(f"[GRAND TALLY] Lock acquired! Pod is aggregating all worker signals for {domain}...")
                                        total_s, total_f = 0, 0
                                        all_sigs = cc_raw.list_blobs(name_starts_with=sig_prefix, include=['metadata'])
                                        for sig in all_sigs:
                                            if sig.name.endswith(".done") and sig.metadata:
                                                total_s += int(sig.metadata.get("s", 0))
                                                total_f += int(sig.metadata.get("f", 0))

                                        end_time = datetime.now(timezone.utc)
                                        raw_start = updated_doc.get("start_time")
                                        start_time = raw_start.replace(tzinfo=timezone.utc) if raw_start and raw_start.tzinfo is None else (raw_start or end_time)
                                        duration_mins = round((end_time - start_time).total_seconds() / 60, 2)

                                        etl_doc = mongo_client["langscope"]["etl_run_history"].find_one({"_id": exec_id})
                                        total_bytes = etl_doc.get("domains", {}).get(domain, {}).get("total_bytes", 0) if etl_doc else 0

                                        idx_col.update_one(
                                            {"_id": exec_id},
                                            {"$set": {
                                                f"domains.{domain}.rows_success": total_s,
                                                f"domains.{domain}.rows_failed": total_f,
                                                f"domains.{domain}.status": "COMPLETED",
                                                f"domains.{domain}.end_time": end_time,
                                                f"domains.{domain}.duration_minutes": duration_mins,
                                                f"domains.{domain}.total_bytes_processed": total_bytes
                                            }}
                                        )
                                        cc_raw.get_blob_client(f"raw/{s_path}/indexing.done").upload_blob(b"DONE", overwrite=True)
                                        logger.info(f"[FINAL SUMMARY] {domain} | Total Success: {total_s:,} | Total Failed: {total_f:,}")
                    except Exception as e: 
                        logger.error(f"[WORKER TALLY CRASH] {e}")

        if not msg_found: time.sleep(5)

if __name__ == "__main__": run_worker()